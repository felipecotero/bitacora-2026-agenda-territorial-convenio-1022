import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colores de la línea gráfica
TINTA = RGBColor(29, 15, 46)      # #1D0F2E
TINTA_2 = RGBColor(58, 42, 82)    # #3A2A52
TINTA_3 = RGBColor(107, 92, 130)  # #6B5C82
CREMA = RGBColor(240, 232, 213)   # #F0E8D5
VIOLETA = RGBColor(90, 54, 128)   # #5A3680
NARANJA = RGBColor(217, 123, 63)  # #D97B3F
VERDE = RGBColor(82, 142, 113)    # #528E71
MAGENTA = RGBColor(168, 67, 138)  # #A8438A
AMARILLO = RGBColor(232, 197, 53) # #E8C535
BLANCO = RGBColor(255, 255, 255)  # #FFFFFF
CARTA = RGBColor(255, 251, 237)   # #FFFBED
LINO = RGBColor(201, 185, 143)    # #C9B98F

def add_slide_base(prs, month_name, month_color):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Fondo Crema
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = CREMA
    
    # Línea inferior (Footer divider)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.5), Inches(11.73), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINO
    line.line.color.rgb = LINO
    
    # Texto Footer
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Agenda Territorial 2026 · Ministerio de las Culturas, las Artes y los Saberes"
    p.font.size = Pt(14)
    p.font.color.rgb = TINTA_3
    p.font.name = "DM Sans"
    
    # Título principal
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Ciclo de Formación Puente"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TINTA
    p.font.name = "Bricolage Grotesque"
    
    # Subtítulo
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.0), Inches(0.6))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Camino Uno — Sesiones de Formación"
    p.font.size = Pt(36)
    p.font.color.rgb = NARANJA
    p.font.name = "Caveat"
    
    # Badge Mes
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.0), Inches(0.6), Inches(2.5), Inches(0.8))
    badge.fill.solid()
    badge.fill.fore_color.rgb = month_color
    badge.line.color.rgb = month_color
    badge.adjustments[0] = 0.5 # Border radius full
    
    tf = badge.text_frame
    tf.vertical_anchor = 3 # Middle
    p = tf.paragraphs[0]
    p.text = month_name
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLANCO
    p.font.name = "Bricolage Grotesque"
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_event(slide, x, y, evt):
    # Fondo de la tarjeta
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.6), Inches(1.0))
    card.fill.solid()
    card.fill.fore_color.rgb = CARTA
    card.line.color.rgb = LINO
    card.adjustments[0] = 0.1
    
    # Colores según orientador
    orient_color = TINTA_3
    text_orient_color = BLANCO
    if evt['orient'] == 'mercy': orient_color = MAGENTA
    elif evt['orient'] == 'jeimmy': orient_color = NARANJA
    elif evt['orient'] == 'felipe': 
        orient_color = AMARILLO
        text_orient_color = TINTA
    elif evt['orient'] == 'yuli': orient_color = VIOLETA
    
    # Cuadro del Date Pill (Mes/Día)
    dp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(y + 0.08), Inches(0.85), Inches(0.85))
    dp.fill.solid()
    dp.fill.fore_color.rgb = orient_color
    dp.line.color.rgb = orient_color
    dp.adjustments[0] = 0.2
    
    # Textos de Date Pill
    dp_tf = dp.text_frame
    dp_tf.word_wrap = False
    dp_tf.margin_left = 0
    dp_tf.margin_right = 0
    dp_tf.margin_top = Inches(0.05)
    
    p = dp_tf.paragraphs[0]
    p.text = f"{evt['mes']}\n"
    p.font.size = Pt(10)
    p.font.color.rgb = text_orient_color
    p.font.bold = True
    p.font.name = "DM Sans"
    p.alignment = PP_ALIGN.CENTER
    
    p2 = dp_tf.add_paragraph()
    p2.text = f"{evt['dia']}\n"
    p2.font.size = Pt(22)
    p2.font.color.rgb = text_orient_color
    p2.font.bold = True
    p2.font.name = "Bricolage Grotesque"
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = dp_tf.add_paragraph()
    p3.text = f"{evt['dow'][:3].upper()}"
    p3.font.size = Pt(9)
    p3.font.color.rgb = text_orient_color
    p3.font.name = "DM Sans"
    p3.alignment = PP_ALIGN.CENTER
    
    # Texto a la derecha
    txt_x = x + 1.1
    
    # Tipo
    txBox = slide.shapes.add_textbox(Inches(txt_x), Inches(y + 0.1), Inches(4.3), Inches(0.3))
    txBox.text_frame.margin_left = 0
    txBox.text_frame.margin_top = 0
    p = txBox.text_frame.paragraphs[0]
    p.text = evt['tipo'].upper()
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = orient_color
    p.font.name = "DM Sans"
    
    # Tema
    txBox = slide.shapes.add_textbox(Inches(txt_x), Inches(y + 0.3), Inches(4.3), Inches(0.4))
    txBox.text_frame.margin_left = 0
    txBox.text_frame.margin_top = 0
    p = txBox.text_frame.paragraphs[0]
    p.text = evt['tema']
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TINTA
    p.font.name = "Bricolage Grotesque"
    
    # Meta (Orientador y hora)
    txBox = slide.shapes.add_textbox(Inches(txt_x), Inches(y + 0.65), Inches(4.3), Inches(0.3))
    txBox.text_frame.margin_left = 0
    txBox.text_frame.margin_top = 0
    p = txBox.text_frame.paragraphs[0]
    
    orient_nombre = {'mercy': 'Mercy', 'jeimmy': 'Jeimmy', 'felipe': 'Felipe', 'yuli': 'Yuli'}.get(evt['orient'], evt['orient'])
    p.text = f"{orient_nombre} · {evt['hora']}"
    p.font.size = Pt(11)
    p.font.color.rgb = TINTA_2
    p.font.name = "DM Sans"

mayo_events = [
    {'mes': 'MAY', 'dia': '11', 'dow': 'Lun', 'orient': 'yuli', 'tipo': 'Sesión inaugural', 'tema': 'Presentación del Ciclo de Formación Puente', 'hora': '18:00–19:30'},
    {'mes': 'MAY', 'dia': '13', 'dow': 'Mié', 'orient': 'mercy', 'tipo': 'Saberes financieros', 'tema': 'Finanzas personales', 'hora': '18:30–20:00'},
    {'mes': 'MAY', 'dia': '15', 'dow': 'Vie', 'orient': 'felipe', 'tipo': 'Comunicaciones', 'tema': 'Desde dónde hablamos — Comunicaciones 1', 'hora': '17:30–19:00'},
    {'mes': 'MAY', 'dia': '19', 'dow': 'Mar', 'orient': 'jeimmy', 'tipo': 'Saberes RIJ', 'tema': 'Contexto: curar el entorno', 'hora': '10:00–11:30'},
    {'mes': 'MAY', 'dia': '22', 'dow': 'Vie', 'orient': 'felipe', 'tipo': 'Comunicaciones', 'tema': 'Decidir qué se dice y qué se cuida — Com. 2', 'hora': '17:30–19:00'},
    {'mes': 'MAY', 'dia': '23', 'dow': 'Sáb', 'orient': 'mercy', 'tipo': 'Saberes financieros', 'tema': 'Diversificación de ingresos', 'hora': '17:30–19:00'},
    {'mes': 'MAY', 'dia': '25', 'dow': 'Lun', 'orient': 'jeimmy', 'tipo': 'Saberes RIJ', 'tema': 'Reciprocidad: la ética de la corresponsabilidad', 'hora': '18:30–20:00'},
]

junio_events = [
    {'mes': 'JUN', 'dia': '01', 'dow': 'Lun', 'orient': 'mercy', 'tipo': 'Saberes financieros', 'tema': 'Flujo de caja', 'hora': '18:30–20:00'},
    {'mes': 'JUN', 'dia': '02', 'dow': 'Mar', 'orient': 'jeimmy', 'tipo': 'Saberes RIJ', 'tema': 'Flujo del ser: la coherencia interna', 'hora': '18:30–20:00'},
    {'mes': 'JUN', 'dia': '03', 'dow': 'Mié', 'orient': 'mercy', 'tipo': 'Saberes financieros', 'tema': 'Indicadores para el control financiero', 'hora': '17:30–19:00'},
    {'mes': 'JUN', 'dia': '05', 'dow': 'Vie', 'orient': 'jeimmy', 'tipo': 'Saberes RIJ', 'tema': 'Herramientas 1: lo político del cuidado situado', 'hora': '17:30–19:00'},
    {'mes': 'JUN', 'dia': '22', 'dow': 'Lun', 'orient': 'jeimmy', 'tipo': 'Saberes RIJ', 'tema': 'Herramientas 2: lo político del cuidado situado', 'hora': '18:30–20:00'},
]

prs = Presentation()
# Set size to Widescreen 16:9 (13.33 x 7.5 inches)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# DIAPOSITIVA MAYO
slide_mayo = add_slide_base(prs, "MAYO", VIOLETA)

# Layout de eventos en mayo (7 eventos: 4 filas x 2 columnas)
start_y = 2.0
gap_y = 1.1
gap_x = 5.9
start_x = 0.8

for idx, evt in enumerate(mayo_events):
    col = idx % 2
    row = idx // 2
    x = start_x + (col * gap_x)
    y = start_y + (row * gap_y)
    add_event(slide_mayo, x, y, evt)

# DIAPOSITIVA JUNIO
slide_junio = add_slide_base(prs, "JUNIO", VERDE)

for idx, evt in enumerate(junio_events):
    col = idx % 2
    row = idx // 2
    x = start_x + (col * gap_x)
    y = start_y + (row * gap_y)
    add_event(slide_junio, x, y, evt)

prs.save('Calendario_Ciclo_Puente.pptx')
print("PPTX generado correctamente: Calendario_Ciclo_Puente.pptx")
